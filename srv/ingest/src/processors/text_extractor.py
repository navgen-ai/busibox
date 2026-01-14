"""
Text extraction from various file formats.

Supports:
- PDF: Marker (primary), TATR (tables), pdfplumber (fallback)
- DOCX: python-docx
- PPTX: python-pptx
- XLSX: openpyxl
- ODT: odfpy
- TXT, HTML, XML, Markdown, CSV, JSON: Direct parsing
- Page image extraction for ColPali (PDFs only)

PDF Splitting:
Large PDFs (>5 pages by default) are automatically split into smaller chunks
before processing to prevent memory issues and timeouts.
"""

import json
import os
import tempfile
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import pdfplumber
import structlog
from docx import Document

from processors.pdf_splitter import PDFSplitter, DEFAULT_PAGES_PER_SPLIT

logger = structlog.get_logger()


def _create_resilient_converter(artifact_dict):
    """
    Create a resilient PDF converter that doesn't fail the entire extraction
    when individual processors (like table recognition) fail.
    
    This patches Marker's PdfConverter.build_document to wrap each processor
    call in try/except, allowing text extraction to succeed even if
    table processing fails.
    """
    from marker.converters.pdf import PdfConverter
    
    # Create a subclass that overrides build_document with resilient processing
    class ResilientPdfConverter(PdfConverter):
        def build_document(self, filepath: str):
            """Build document with resilient processor handling."""
            from marker.providers.registry import provider_from_filepath
            from marker.builders.document import DocumentBuilder
            from marker.builders.structure import StructureBuilder
            from marker.builders.line import LineBuilder
            from marker.builders.ocr import OcrBuilder
            
            provider_cls = provider_from_filepath(filepath)
            layout_builder = self.resolve_dependencies(self.layout_builder_class)
            line_builder = self.resolve_dependencies(LineBuilder)
            ocr_builder = self.resolve_dependencies(OcrBuilder)
            provider = provider_cls(filepath, self.config)
            document = DocumentBuilder(self.config)(
                provider, layout_builder, line_builder, ocr_builder
            )
            structure_builder_cls = self.resolve_dependencies(StructureBuilder)
            structure_builder_cls(document)

            # Run processors with individual error handling
            failed_processors = []
            for processor in self.processor_list:
                processor_name = processor.__class__.__name__
                try:
                    processor(document)
                except Exception as e:
                    # Log the failure but continue with other processors
                    logger.warning(
                        "Marker processor failed, continuing without it",
                        processor=processor_name,
                        error=str(e),
                        error_type=type(e).__name__,
                    )
                    failed_processors.append(processor_name)
            
            if failed_processors:
                logger.info(
                    "Marker completed with some processor failures",
                    failed_processors=failed_processors,
                    total_processors=len(self.processor_list),
                )

            return document
    
    return ResilientPdfConverter(artifact_dict=artifact_dict)


class ExtractionResult:
    """Result of text extraction."""
    
    def __init__(
        self,
        text: str,
        markdown: Optional[str] = None,
        page_images: Optional[List[str]] = None,
        page_count: int = 0,
        tables: Optional[List[Dict]] = None,
        metadata: Optional[Dict] = None,
    ):
        self.text = text
        self.markdown = markdown
        self.page_images = page_images or []
        self.page_count = page_count
        self.tables = tables or []
        self.metadata = metadata or {}


class TextExtractor:
    """Extract text from various file formats."""
    
    # Class-level cache for Marker models to avoid reloading on every extraction
    # This prevents OOM when processing multiple documents
    _marker_models = None
    _marker_converter = None
    
    @classmethod
    def get_marker_models(cls):
        """Get or create cached Marker models (singleton pattern).
        
        Uses ResilientPdfConverter which wraps individual processors in
        try/except so that table recognition failures don't cause the
        entire extraction to fall back to pdfplumber.
        """
        if cls._marker_models is None:
            try:
                from marker.models import create_model_dict
                
                logger.info("Loading Marker models (will be cached for reuse)...")
                cls._marker_models = create_model_dict()
                # Use resilient converter that handles processor failures gracefully
                cls._marker_converter = _create_resilient_converter(cls._marker_models)
                logger.info("Marker models loaded and cached (resilient mode)", models=list(cls._marker_models.keys()))
            except Exception as e:
                logger.error("Failed to load Marker models", error=str(e))
                raise
        return cls._marker_models, cls._marker_converter
    
    @classmethod
    def cleanup_marker_models(cls):
        """Clean up cached Marker models to free GPU memory."""
        if cls._marker_models is not None:
            logger.info("Cleaning up Marker models...")
            del cls._marker_models
            del cls._marker_converter
            cls._marker_models = None
            cls._marker_converter = None
            
            # Force GPU memory cleanup
            try:
                import torch
                if torch.cuda.is_available():
                    torch.cuda.empty_cache()
                    logger.info("GPU memory cleared")
            except Exception:
                pass
    
    def __init__(self, config: dict):
        """Initialize text extractor."""
        self.config = config
        self.temp_dir = config.get("temp_dir", "/tmp/ingest")
        self.marker_enabled = config.get("marker_enabled", True)  # Can disable to save memory
        
        # PDF splitting configuration
        # Split large PDFs into chunks of this many pages (default: 5)
        self.pdf_split_pages = config.get("pdf_split_pages", DEFAULT_PAGES_PER_SPLIT)
        # Enable/disable PDF splitting (default: True)
        self.pdf_split_enabled = config.get("pdf_split_enabled", True)
        
        # Initialize PDF splitter
        self.pdf_splitter = PDFSplitter(
            pages_per_split=self.pdf_split_pages,
            temp_dir=self.temp_dir,
        )
        
        # Remote Marker service URL - if set, calls remote service instead of local Marker
        # This allows test environment to use production Marker service
        self.marker_service_url = config.get("marker_service_url") or os.getenv("MARKER_SERVICE_URL")
        if self.marker_service_url:
            logger.info(
                "Using remote Marker service",
                url=self.marker_service_url,
            )
        
        if self.pdf_split_enabled:
            logger.info(
                "PDF splitting enabled",
                pages_per_split=self.pdf_split_pages,
            )
        
        os.makedirs(self.temp_dir, exist_ok=True)
    
    def extract(self, file_path: str, mime_type: str) -> ExtractionResult:
        """
        Extract text from file.
        
        Args:
            file_path: Path to file
            mime_type: MIME type of file
            
        Returns:
            ExtractionResult with text, markdown, page images, etc.
        """
        # PDF
        if mime_type == "application/pdf":
            return self._extract_pdf(file_path)
        
        # Microsoft Office formats
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return self._extract_docx(file_path)
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return self._extract_pptx(file_path)
        elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return self._extract_xlsx(file_path)
        
        # OpenDocument formats
        elif mime_type == "application/vnd.oasis.opendocument.text":
            return self._extract_odt(file_path)
        elif mime_type == "application/vnd.oasis.opendocument.presentation":
            return self._extract_odp(file_path)
        elif mime_type == "application/vnd.oasis.opendocument.spreadsheet":
            return self._extract_ods(file_path)
        
        # Text formats
        elif mime_type == "text/plain":
            return self._extract_txt(file_path)
        elif mime_type == "text/html":
            return self._extract_html(file_path)
        elif mime_type in ("text/xml", "application/xml"):
            return self._extract_xml(file_path)
        elif mime_type == "text/markdown":
            return self._extract_markdown(file_path)
        
        # Data formats
        elif mime_type == "text/csv":
            return self._extract_csv(file_path)
        elif mime_type == "application/json":
            return self._extract_json(file_path)
        
        else:
            raise ValueError(f"Unsupported MIME type: {mime_type}")
    
    def _extract_pdf(self, file_path: str) -> ExtractionResult:
        """Extract text from PDF using Marker, TATR, and page images.
        
        For large PDFs (>pdf_split_pages pages), the PDF is split into smaller
        chunks before processing to prevent memory issues and timeouts.
        Results from each chunk are then combined.
        """
        # Check if PDF needs splitting
        if self.pdf_split_enabled and self.pdf_splitter.needs_splitting(file_path):
            return self._extract_pdf_with_splitting(file_path)
        
        # Process as single PDF (no splitting needed)
        return self._extract_single_pdf(file_path)
    
    def _extract_pdf_with_splitting(self, file_path: str) -> ExtractionResult:
        """Extract text from a large PDF by splitting it into chunks.
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Combined ExtractionResult from all chunks
        """
        total_page_count = self.pdf_splitter.get_page_count(file_path)
        
        logger.info(
            "Processing large PDF with splitting",
            file_path=file_path,
            total_pages=total_page_count,
            pages_per_split=self.pdf_split_pages,
        )
        
        # Split the PDF
        splits = self.pdf_splitter.split(file_path)
        
        # Accumulate results from all splits
        all_text_parts = []
        all_markdown_parts = []
        all_page_images = []
        all_tables = []
        extraction_method = "unknown"
        
        try:
            for split_idx, (split_path, start_page, end_page) in enumerate(splits):
                logger.info(
                    "Processing PDF split",
                    split_num=split_idx + 1,
                    total_splits=len(splits),
                    pages=f"{start_page}-{end_page}",
                    split_path=split_path,
                )
                
                # Extract from this split
                logger.debug(f"Calling _extract_single_pdf for split {split_idx + 1}")
                result = self._extract_single_pdf(split_path)
                logger.debug(f"_extract_single_pdf returned for split {split_idx + 1}")
                
                # Accumulate text
                logger.debug(f"Accumulating text for split {split_idx + 1}")
                if result.text:
                    # Add page range marker for debugging/reference
                    all_text_parts.append(result.text)
                logger.debug(f"Text accumulated for split {split_idx + 1}")
                
                # Accumulate markdown
                logger.debug(f"Accumulating markdown for split {split_idx + 1}")
                if result.markdown:
                    all_markdown_parts.append(result.markdown)
                logger.debug(f"Markdown accumulated for split {split_idx + 1}")
                
                # Accumulate page images (adjust paths to avoid conflicts)
                logger.debug(f"Accumulating page images for split {split_idx + 1}")
                if result.page_images:
                    all_page_images.extend(result.page_images)
                logger.debug(f"Page images accumulated for split {split_idx + 1}: {len(result.page_images) if result.page_images else 0}")
                
                # Accumulate tables
                logger.debug(f"Accumulating tables for split {split_idx + 1}")
                if result.tables:
                    # Add page offset to table metadata
                    for table in result.tables:
                        if isinstance(table, dict):
                            table["page_offset"] = start_page - 1
                    all_tables.extend(result.tables)
                logger.debug(f"Tables accumulated for split {split_idx + 1}")
                
                # Use extraction method from first successful extraction
                logger.debug(f"Checking extraction method for split {split_idx + 1}")
                if extraction_method == "unknown" and result.metadata.get("extraction_method"):
                    extraction_method = result.metadata.get("extraction_method")
                logger.debug(f"Extraction method set: {extraction_method}")
                
                logger.info(
                    "Split extraction complete",
                    split_num=split_idx + 1,
                    text_length=len(result.text) if result.text else 0,
                    page_images=len(result.page_images) if result.page_images else 0,
                    has_markdown=result.markdown is not None,
                )
        
        finally:
            # Clean up split files
            self.pdf_splitter.cleanup_splits(splits, file_path)
        
        # Combine results
        combined_text = "\n\n".join(all_text_parts)
        combined_markdown = "\n\n".join(all_markdown_parts) if all_markdown_parts else None
        
        logger.info(
            "PDF splitting extraction complete",
            file_path=file_path,
            total_pages=total_page_count,
            num_splits=len(splits),
            combined_text_length=len(combined_text),
            combined_page_images=len(all_page_images),
            combined_tables=len(all_tables),
        )
        
        return ExtractionResult(
            text=combined_text,
            markdown=combined_markdown,
            page_images=all_page_images,
            page_count=total_page_count,
            tables=all_tables,
            metadata={
                "extraction_method": extraction_method,
                "split_processing": True,
                "num_splits": len(splits),
                "pages_per_split": self.pdf_split_pages,
            },
        )
    
    def _extract_single_pdf(self, file_path: str) -> ExtractionResult:
        """Extract text from a single PDF (no splitting).
        
        This is the original extraction logic, now used for both small PDFs
        and individual chunks of split PDFs.
        """
        page_images = []
        markdown_text = None
        text_content = ""
        tables = []
        page_count = 0
        extraction_method = "unknown"
        
        # Check if Marker is enabled (can be disabled to save memory)
        if not self.marker_enabled:
            logger.info("Marker disabled, using pdfplumber fallback", file_path=file_path)
            markdown_text = None  # Skip Marker, go straight to pdfplumber
        elif self.marker_service_url:
            # Use remote Marker service
            markdown_text, page_count, extraction_method = self._extract_pdf_remote_marker(file_path)
            if markdown_text:
                text_content = markdown_text
                # Extract page images locally for ColPali
                page_images = self._extract_pdf_page_images(file_path)
        else:
            # Try Marker first (best quality)
            # Marker v1.x uses a different API than v0.x
            try:
                logger.debug("Attempting to import marker (trying v1.x API first)")
                try:
                    # Try marker v1.x API (marker-pdf >= 1.0)
                    from marker.converters.pdf import PdfConverter
                    from marker.models import create_model_dict
                    import torch
                    import os
                    
                    # Configure GPU settings for Marker
                    marker_use_gpu = self.config.get("marker_use_gpu", True)
                    marker_gpu_device = self.config.get("marker_gpu_device", "cuda")
                    
                    # Determine device
                    if marker_use_gpu and torch.cuda.is_available():
                        device = marker_gpu_device if marker_gpu_device != "auto" else "cuda"
                        logger.info(
                            "Marker GPU available and enabled",
                            file_path=file_path,
                            device=device,
                            cuda_device_count=torch.cuda.device_count(),
                        )
                    else:
                        device = "cpu"
                        if marker_use_gpu and not torch.cuda.is_available():
                            logger.warning(
                                "Marker GPU requested but CUDA not available, using CPU",
                                file_path=file_path,
                            )
                        else:
                            logger.info("Marker using CPU", file_path=file_path)
                    
                    # Set environment variables for Marker
                    os.environ["TORCH_DEVICE"] = device
                    if device == "cuda":
                        # Set GPU memory settings
                        inference_ram = self.config.get("marker_inference_ram", "16")
                        vram_per_task = self.config.get("marker_vram_per_task", "3.5")
                        os.environ["INFERENCE_RAM"] = str(inference_ram)
                        os.environ["VRAM_PER_TASK"] = str(vram_per_task)
                        logger.debug(
                            "Marker GPU memory configured",
                            inference_ram=inference_ram,
                            vram_per_task=vram_per_task,
                        )
                    
                    logger.info("Using Marker v1.x for PDF extraction", file_path=file_path, device=device)
                    
                    # Use cached models to avoid OOM on multiple extractions
                    # Models are loaded once and reused across all extractions
                    artifact_dict, converter = self.get_marker_models()
                    
                    # Convert PDF to markdown
                    result = converter(file_path)
                    markdown_text = result.markdown
                    text_content = markdown_text
                    
                    # Extract page images for ColPali
                    page_images = self._extract_pdf_page_images(file_path)
                    page_count = len(page_images)
                    
                except ImportError:
                    # Fall back to marker v0.x API (old marker-pdf)
                    logger.debug("Marker v1.x not found, trying v0.x API")
                    from marker.convert import convert_single_pdf
                    
                    logger.info("Using Marker v0.x for PDF extraction", file_path=file_path)
                    markdown_text, images, metadata = convert_single_pdf(file_path)
                    text_content = markdown_text
                    page_count = len(images) if images else 0
                    
                    # Extract page images for ColPali
                    page_images = self._extract_pdf_page_images(file_path)
                    page_count = len(page_images)
                
            except ImportError as e:
                logger.warning(
                    "Marker not available, falling back to pdfplumber",
                    error=str(e),
                    error_type=type(e).__name__,
                    import_error_details=getattr(e, 'name', 'unknown'),
                )
                markdown_text = None
            except Exception as e:
                # Log as warning since we'll fall back to pdfplumber
                # Common issues: table_rec tensor shape errors, memory issues
                logger.warning(
                    "Marker execution failed, falling back to pdfplumber",
                    error=str(e),
                    error_type=type(e).__name__,
                )
                markdown_text = None
            
            # Extract tables with TATR if available
            try:
                from tatr import TableTransformer
                
                logger.info("Extracting tables with TATR", file_path=file_path)
                table_transformer = TableTransformer()
                tables = table_transformer.extract_tables(file_path)
                
                # Add table text to content
                for table in tables:
                    if isinstance(table, dict) and "text" in table:
                        text_content += "\n\n" + table["text"]
            except ImportError:
                logger.debug("TATR not available, skipping table extraction")
        
        # Fallback to pdfplumber if Marker failed or disabled
        if not text_content or page_count == 0:
            logger.info("Using pdfplumber fallback", file_path=file_path)
            try:
                with pdfplumber.open(file_path) as pdf:
                    page_count = len(pdf.pages)
                    text_parts = []
                    
                    for page in pdf.pages:
                        try:
                            # Default extraction works best - layout=True adds too much whitespace
                            page_text = page.extract_text()
                            if page_text:
                                text_parts.append(page_text)
                        except (IndexError, KeyError) as e:
                            # Some malformed PDFs cause pdfplumber to fail on specific pages
                            logger.warning(
                                "pdfplumber failed on page, skipping",
                                page_number=page.page_number,
                                error=str(e),
                                error_type=type(e).__name__,
                            )
                            continue
                    
                    text_content = "\n\n".join(text_parts)  # Double newline between pages
                    
                    # Extract page images if not already done
                    if not page_images:
                        page_images = self._extract_pdf_page_images(file_path)
            except Exception as e:
                logger.error(
                    "pdfplumber extraction failed completely",
                    error=str(e),
                    error_type=type(e).__name__,
                    file_path=file_path,
                )
                # Continue with empty text - will trigger OCR or other fallback
                text_content = ""
                page_count = 0
        
        # Detect scanned PDFs (no extractable text)
        logger.debug(
            "Checking if OCR needed",
            file_path=file_path,
            has_text=bool(text_content.strip()),
            page_count=page_count,
        )
        if not text_content.strip() and page_count > 0:
            logger.warning("No text extracted, PDF may be scanned - OCR required", file_path=file_path)
            # Trigger OCR processing (implemented separately)
            text_content = self._ocr_pdf(file_path)
        
        # Determine extraction method for metadata
        if not extraction_method or extraction_method == "unknown":
            extraction_method = "marker" if markdown_text else "pdfplumber"
        
        # Clear GPU memory after extraction to prevent accumulation
        try:
            import torch
            if torch.cuda.is_available():
                torch.cuda.empty_cache()
                logger.debug("Cleared CUDA cache after PDF extraction")
        except Exception as e:
            logger.debug("GPU cache clear skipped", error=str(e))
        
        logger.info(
            "PDF extraction complete, returning result",
            file_path=file_path,
            text_length=len(text_content),
            page_count=page_count,
            has_page_images=len(page_images) > 0,
        )
        
        return ExtractionResult(
            text=text_content,
            markdown=markdown_text,
            page_images=page_images,
            page_count=page_count,
            tables=tables,
            metadata={"extraction_method": extraction_method},
        )
    
    def _extract_pdf_remote_marker(self, file_path: str) -> Tuple[Optional[str], int, str]:
        """
        Extract text from PDF using remote Marker service.
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Tuple of (markdown_text, page_count, extraction_method)
        """
        import httpx
        
        logger.info(
            "Calling remote Marker service",
            file_path=file_path,
            service_url=self.marker_service_url,
        )
        
        try:
            with open(file_path, "rb") as f:
                files = {"file": (Path(file_path).name, f, "application/pdf")}
                
                with httpx.Client(timeout=300.0) as client:  # 5 minute timeout for large PDFs
                    response = client.post(
                        f"{self.marker_service_url}/extract",
                        files=files,
                    )
                    
                    if response.status_code != 200:
                        logger.warning(
                            "Remote Marker service failed",
                            status=response.status_code,
                            response=response.text[:200],
                        )
                        return None, 0, "remote_marker_failed"
                    
                    result = response.json()
                    markdown_text = result.get("text", "") or result.get("markdown", "")
                    page_count = result.get("page_count", 0)
                    
                    logger.info(
                        "Remote Marker extraction complete",
                        text_length=len(markdown_text),
                        page_count=page_count,
                    )
                    
                    return markdown_text, page_count, "remote_marker"
                    
        except Exception as e:
            logger.error(
                "Remote Marker service error",
                error=str(e),
                error_type=type(e).__name__,
            )
            return None, 0, "remote_marker_error"
    
    def _extract_pdf_page_images(self, file_path: str) -> List[str]:
        """Extract page images from PDF for ColPali.
        
        Images are scaled to ensure ColPali generates at most 32 patches (4096 dims)
        to avoid truncation and information loss.
        
        Returns empty list if poppler-utils is not installed or extraction fails.
        """
        try:
            from pdf2image import convert_from_path
            from pdf2image.exceptions import (
                PDFInfoNotInstalledError,
                PDFPageCountError,
                PDFSyntaxError
            )
            from PIL import Image
            
            logger.info("Extracting PDF page images", file_path=file_path)
            
            # pdf2image needs poppler-utils (pdftoppm) installed
            try:
                images = convert_from_path(file_path, dpi=150)
            except (PDFInfoNotInstalledError, PDFPageCountError) as e:
                # Poppler not installed or can't read PDF
                logger.warning(
                    "Poppler not available, skipping page image extraction",
                    file_path=file_path,
                    error=str(e),
                    error_type=type(e).__name__,
                )
                return []
            except PDFSyntaxError as e:
                # Malformed PDF
                logger.warning(
                    "PDF syntax error, skipping page image extraction",
                    file_path=file_path,
                    error=str(e),
                )
                return []
            except Exception as conv_error:
                # Other conversion errors
                logger.error(
                    "Page image extraction failed", 
                    file_path=file_path, 
                    error=str(conv_error),
                    error_type=type(conv_error).__name__,
                )
                return []
            
            if not images:
                logger.warning("No page images extracted", file_path=file_path)
                return []
            
            # ColPali target: 32 patches max (32 * 128 = 4096 dims)
            # Each patch is roughly 14x14 pixels, so 32 patches ≈ 448x448 pixels
            # Scale images to max 960x960 to stay under 32 patches
            max_dimension = 960
            
            # Save images to temp directory
            page_images = []
            base_name = Path(file_path).stem
            
            for i, image in enumerate(images):
                # Scale down large images to prevent truncation
                width, height = image.size
                if width > max_dimension or height > max_dimension:
                    # Calculate scaling factor to fit within max_dimension
                    scale = max_dimension / max(width, height)
                    new_width = int(width * scale)
                    new_height = int(height * scale)
                    image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)
                    logger.debug(
                        "Scaled page image",
                        page=i+1,
                        original_size=f"{width}x{height}",
                        scaled_size=f"{new_width}x{new_height}",
                    )
                
                image_path = os.path.join(self.temp_dir, f"{base_name}_page_{i+1:03d}.png")
                image.save(image_path, "PNG")
                page_images.append(image_path)
            
            logger.info("Extracted page images", count=len(page_images), file_path=file_path)
            return page_images
        
        except ImportError as e:
            # pdf2image or pdf2image.exceptions not available
            logger.warning(
                "pdf2image not available, skipping page image extraction",
                error=str(e),
            )
            return []
        except Exception as e:
            logger.error(
                "Page image extraction failed unexpectedly", 
                file_path=file_path, 
                error=str(e),
                error_type=type(e).__name__,
            )
            return []
    
    def _ocr_pdf(self, file_path: str) -> str:
        """Perform OCR on scanned PDF."""
        try:
            from pdf2image import convert_from_path
            import pytesseract
            
            logger.info("Performing OCR on scanned PDF", file_path=file_path)
            images = convert_from_path(file_path, dpi=300)
            
            text_parts = []
            for image in images:
                text = pytesseract.image_to_string(image)
                text_parts.append(text)
            
            return "\n".join(text_parts)
        
        except ImportError:
            logger.warning("OCR dependencies not available (pdf2image, pytesseract)")
            return ""
        except Exception as e:
            logger.error("OCR failed", file_path=file_path, error=str(e))
            return ""
    
    def _extract_docx(self, file_path: str) -> ExtractionResult:
        """Extract text from DOCX file."""
        try:
            doc = Document(file_path)
            text_parts = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append({"data": table_data})
            
            text_content = "\n".join(text_parts)
            
            return ExtractionResult(
                text=text_content,
                page_count=len(doc.paragraphs) // 20,  # Estimate pages
                tables=tables,
                metadata={"extraction_method": "python-docx"},
            )
        
        except Exception as e:
            logger.error("DOCX extraction failed", file_path=file_path, error=str(e), exc_info=True)
            raise
    
    def _extract_txt(self, file_path: str) -> ExtractionResult:
        """Extract text from TXT file."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            
            return ExtractionResult(
                text=text,
                page_count=len(text.split("\n")) // 50,  # Estimate pages
                metadata={"extraction_method": "direct_read"},
            )
        
        except Exception as e:
            logger.error("TXT extraction failed", file_path=file_path, error=str(e), exc_info=True)
            raise
    
    def _extract_html(self, file_path: str) -> ExtractionResult:
        """Extract text from HTML file."""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, "r", encoding="utf-8") as f:
                html = f.read()
            
            soup = BeautifulSoup(html, "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            
            return ExtractionResult(
                text=text,
                page_count=1,
                metadata={"extraction_method": "beautifulsoup"},
            )
        
        except ImportError:
            # Fallback: basic regex extraction
            import re
            with open(file_path, "r", encoding="utf-8") as f:
                html = f.read()
            text = re.sub(r"<[^>]+>", "", html)
            
            return ExtractionResult(
                text=text,
                page_count=1,
                metadata={"extraction_method": "regex_fallback"},
            )
    
    def _extract_markdown(self, file_path: str) -> ExtractionResult:
        """Extract text from Markdown file."""
        return self._extract_txt(file_path)  # Markdown is plain text
    
    def _extract_csv(self, file_path: str) -> ExtractionResult:
        """Extract text from CSV file."""
        try:
            import csv
            
            text_parts = []
            with open(file_path, "r", encoding="utf-8") as f:
                reader = csv.reader(f)
                for row in reader:
                    text_parts.append(" | ".join(row))
            
            text = "\n".join(text_parts)
            
            return ExtractionResult(
                text=text,
                page_count=1,
                metadata={"extraction_method": "csv"},
            )
        
        except Exception as e:
            logger.error("CSV extraction failed", file_path=file_path, error=str(e), exc_info=True)
            raise
    
    def _extract_json(self, file_path: str) -> ExtractionResult:
        """Extract text from JSON file."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            
            # Convert JSON to readable text
            text = json.dumps(data, indent=2)
            
            return ExtractionResult(
                text=text,
                page_count=1,
                metadata={"extraction_method": "json"},
            )
        
        except Exception as e:
            logger.error("JSON extraction failed", file_path=file_path, error=str(e), exc_info=True)
            raise
    
    def _extract_xml(self, file_path: str) -> ExtractionResult:
        """Extract text from XML file."""
        try:
            from lxml import etree
            
            # Parse XML
            tree = etree.parse(file_path)
            root = tree.getroot()
            
            # Extract all text content
            text_parts = []
            for element in root.iter():
                if element.text and element.text.strip():
                    text_parts.append(element.text.strip())
                if element.tail and element.tail.strip():
                    text_parts.append(element.tail.strip())
            
            text = "\n".join(text_parts)
            
            return ExtractionResult(
                text=text,
                page_count=1,
                metadata={"extraction_method": "lxml"},
            )
        
        except ImportError:
            # Fallback to built-in xml.etree
            import xml.etree.ElementTree as ET
            
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            text_parts = []
            for element in root.iter():
                if element.text and element.text.strip():
                    text_parts.append(element.text.strip())
            
            text = "\n".join(text_parts)
            
            return ExtractionResult(
                text=text,
                page_count=1,
                metadata={"extraction_method": "xml.etree"},
            )
        
        except Exception as e:
            logger.error("XML extraction failed", file_path=file_path, error=str(e), exc_info=True)
            raise
    
    def _extract_pptx(self, file_path: str) -> ExtractionResult:
        """Extract text from PowerPoint (PPTX) file."""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_parts = []
            slide_count = 0
            
            for slide in prs.slides:
                slide_count += 1
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text_parts.append(f"=== Slide {slide_count} ===")
                    text_parts.extend(slide_text)
                    text_parts.append("")  # Empty line between slides
            
            text = "\n".join(text_parts)
            
            return ExtractionResult(
                text=text,
                page_count=slide_count,
                metadata={
                    "extraction_method": "python-pptx",
                    "slide_count": slide_count,
                },
            )
        
        except ImportError:
            logger.error("python-pptx not installed", file_path=file_path)
            raise ValueError("python-pptx library required for PPTX extraction")
        except Exception as e:
            logger.error("PPTX extraction failed", file_path=file_path, error=str(e), exc_info=True)
            raise
    
    def _extract_xlsx(self, file_path: str) -> ExtractionResult:
        """Extract text from Excel (XLSX) file."""
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            sheet_count = 0
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_count += 1
                
                text_parts.append(f"=== Sheet: {sheet_name} ===")
                
                # Extract cell values row by row
                for row in sheet.iter_rows(values_only=True):
                    # Filter out empty cells and convert to strings
                    row_values = [str(cell) for cell in row if cell is not None]
                    if row_values:
                        text_parts.append(" | ".join(row_values))
                
                text_parts.append("")  # Empty line between sheets
            
            wb.close()
            text = "\n".join(text_parts)
            
            return ExtractionResult(
                text=text,
                page_count=sheet_count,
                metadata={
                    "extraction_method": "openpyxl",
                    "sheet_count": sheet_count,
                },
            )
        
        except ImportError:
            logger.error("openpyxl not installed", file_path=file_path)
            raise ValueError("openpyxl library required for XLSX extraction")
        except Exception as e:
            logger.error("XLSX extraction failed", file_path=file_path, error=str(e), exc_info=True)
            raise
    
    def _extract_odt(self, file_path: str) -> ExtractionResult:
        """Extract text from OpenDocument Text (ODT) file."""
        try:
            from odf import text, teletype
            from odf.opendocument import load
            
            doc = load(file_path)
            text_parts = []
            
            # Extract all text elements
            for paragraph in doc.getElementsByType(text.P):
                para_text = teletype.extractText(paragraph)
                if para_text.strip():
                    text_parts.append(para_text)
            
            # Extract text from tables
            for table in doc.getElementsByType(text.Table):
                for row in table.getElementsByType(text.TableRow):
                    row_text = []
                    for cell in row.getElementsByType(text.TableCell):
                        cell_text = teletype.extractText(cell).strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            text_content = "\n".join(text_parts)
            
            return ExtractionResult(
                text=text_content,
                page_count=len(text_parts) // 20,  # Estimate pages
                metadata={"extraction_method": "odfpy"},
            )
        
        except ImportError:
            logger.error("odfpy not installed", file_path=file_path)
            raise ValueError("odfpy library required for ODT extraction")
        except Exception as e:
            logger.error("ODT extraction failed", file_path=file_path, error=str(e), exc_info=True)
            raise
    
    def _extract_odp(self, file_path: str) -> ExtractionResult:
        """Extract text from OpenDocument Presentation (ODP) file."""
        try:
            from odf import text, teletype
            from odf.opendocument import load
            
            doc = load(file_path)
            text_parts = []
            slide_count = 0
            
            # ODP presentations have draw:page elements
            from odf.draw import Page
            for page in doc.getElementsByType(Page):
                slide_count += 1
                slide_text = []
                
                # Extract text from all text elements on the slide
                for paragraph in page.getElementsByType(text.P):
                    para_text = teletype.extractText(paragraph)
                    if para_text.strip():
                        slide_text.append(para_text)
                
                if slide_text:
                    text_parts.append(f"=== Slide {slide_count} ===")
                    text_parts.extend(slide_text)
                    text_parts.append("")
            
            text_content = "\n".join(text_parts)
            
            return ExtractionResult(
                text=text_content,
                page_count=slide_count,
                metadata={
                    "extraction_method": "odfpy",
                    "slide_count": slide_count,
                },
            )
        
        except ImportError:
            logger.error("odfpy not installed", file_path=file_path)
            raise ValueError("odfpy library required for ODP extraction")
        except Exception as e:
            logger.error("ODP extraction failed", file_path=file_path, error=str(e), exc_info=True)
            raise
    
    def _extract_ods(self, file_path: str) -> ExtractionResult:
        """Extract text from OpenDocument Spreadsheet (ODS) file."""
        try:
            from odf import table, teletype
            from odf.opendocument import load
            
            doc = load(file_path)
            text_parts = []
            sheet_count = 0
            
            # Extract text from all tables (sheets)
            for spreadsheet_table in doc.getElementsByType(table.Table):
                sheet_count += 1
                sheet_name = spreadsheet_table.getAttribute("name")
                text_parts.append(f"=== Sheet: {sheet_name} ===")
                
                # Extract rows
                for row in spreadsheet_table.getElementsByType(table.TableRow):
                    row_text = []
                    for cell in row.getElementsByType(table.TableCell):
                        cell_text = teletype.extractText(cell).strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        text_parts.append(" | ".join(row_text))
                
                text_parts.append("")
            
            text_content = "\n".join(text_parts)
            
            return ExtractionResult(
                text=text_content,
                page_count=sheet_count,
                metadata={
                    "extraction_method": "odfpy",
                    "sheet_count": sheet_count,
                },
            )
        
        except ImportError:
            logger.error("odfpy not installed", file_path=file_path)
            raise ValueError("odfpy library required for ODS extraction")
        except Exception as e:
            logger.error("ODS extraction failed", file_path=file_path, error=str(e), exc_info=True)
            raise
